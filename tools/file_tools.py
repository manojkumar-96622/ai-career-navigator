
import base64
from pptx import Presentation
from docx import Document
from pypdf import PdfReader
from fpdf import FPDF

def read_ppt_file_from_stream(uploaded_file):
    try:
        prs = Presentation(uploaded_file)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"): text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e: return f"Error extracting PPTX: {str(e)}"

def read_docx_file_from_stream(uploaded_file):
    try:
        doc = Document(uploaded_file)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e: return f"Error extracting DOCX: {str(e)}"

def read_pdf_file_from_stream(uploaded_file):
    try:
        reader = PdfReader(uploaded_file)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e: return f"Error extracting PDF: {str(e)}"

def convert_to_pdf(content):
    try:
        # Normalize whitespace and remove problematic characters
        content = content.replace('\x00', '').replace('\r', '').replace('\t', '    ')
        replacements = {
            '\u2013': '-', '\u2014': '-', '\u2018': "'", '\u2019': "'", 
            '\u201c': '"', '\u201d': '"', '\u2022': '*', '\u2026': '...'
        }
        for k, v in replacements.items(): content = content.replace(k, v)

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_margins(left=15, top=15, right=15)
        pdf.set_font("Helvetica", size=11)
        
        # Render line by line with explicit width (A4 210mm - 15mm*2 margins = 180mm)
        for line in content.split('\n'):
            safe_line = line.encode('latin-1', 'replace').decode('latin-1')
            pdf.multi_cell(180, 6, txt=safe_line)
        
        pdf_output = pdf.output()
        b64 = base64.b64encode(pdf_output).decode('utf-8')
        return {"pdf_url": f"data:application/pdf;base64,{b64}"}
    except Exception as e:
        import logging
        logging.error(f"PDF Tool Failure: {str(e)}")
        return {"pdf_url": None, "error": f"Failed to generate PDF: {str(e)}"}
